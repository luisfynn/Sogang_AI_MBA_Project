import collections
import collections.abc

####1. 파워포인트 문서 작성
from pptx import Presentation
from pptx.util import Inches
prs = Presentation()

for i in range(0,11):
    title_slide_layout = prs.slide_layouts[i]    # 슬라이드 종류 선택
    slide = prs.slides.add_slide(title_slide_layout) # 슬라이드 추가
prs.save('add all slides.pptx') # 내용 확인한 후 닫기(경로가 지정되지 않으면 기본 경로에 저장됨)

# [2] 레이아웃 별 placeholder 속성 확인
for i in range(0,11):
    print('----[%d]----'%(i))
    slide = prs.slides.add_slide(prs.slide_layouts[i])
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))
# 상세 내용 보기
slide4 = prs.slides.add_slide(prs.slide_layouts[4]) # 슬라이드4 placeholder
slide4.placeholders[0].placeholder_format.idx   # 슬라이드 4의 첫번째 index
slide4.placeholders[0].name     #슬라이드4의 첫번째 속성 (큰 제목)

#파이썬으로 파워포인트 작성(타이틀 페이지)
from pptx import Presentation
from pptx.util import Inches
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]    # 제목과 부제목만 있는 레이아웃 선택
slide = prs.slides.add_slide(title_slide_layout) # 슬라이드 추가
# 제목 작성
title = slide.placeholders[0]
title.text = '문서 자동화 로봇'
# 부제목 작성
subtitle = slide.placeholders[1]
subtitle.text = 'python pptx 자동화 프로그램'
# 저장
prs.save('test1.pptx')
# 내용 확인 후 pptx 화면 닫기

#파이썬으로 파워포인트 작성(본문페이지: 타이틀과 내용)
bullet_slide_layout = prs.slide_layouts[1] # 제목과 bullet 으로 내용을 입력하는 레이아웃 선택
slide = prs.slides.add_slide(bullet_slide_layout) # 기존 슬라이드에 추가
#제목 작성
title_shape = slide.placeholders[0]
title_shape.text = 'Bullet 추가'
#내용
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = 'bullet slide layout 찾기'
#단락 추가
p = tf.add_paragraph()
p.text = 'TextFrame.txt 사용하기'
p.level = 1 # 들여쓰기 레벨
#저장
prs.save('test1.pptx')
#내용 확인 후 pptx 화면 닫기

#파이썬으로 파워포인트 작성(본문 페이지: 사진 삽입)
img_path = 'cat.jpg'         # 원본 사진이 저장된 경로
# 무료 사진/동영상 다운로드 https://www.pexels.com/ko-kr/
blank_slide_layout = prs.slide_layouts[6]      # 아무런 포맷이 없는 레이아웃 선택
slide = prs.slides.add_slide(blank_slide_layout)
# 작게 편집한 사진 삽입
left = top = Inches(1)
width = height = Inches(1)
# width, height가 없을 경우 원본 사이즈 이용
pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)
# 크게 편집한 사진 삽입
left = Inches(3)
width = Inches(5.5)
height = Inches(4)
pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)
prs.save('test1.pptx')

#파이썬으로 파워포인트 작성(본문 페이지: 표 삽입)
title_only_slide_layout = prs.slide_layouts[5]    # 제목만 있는 레이아웃 선택
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes
title_shape = slide.placeholders[0]
title_shape.text = 'Table 1. 회사명과 주식가격'
rows = 3
cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)
table = shapes.add_table(rows, cols, left, top, width, height).table
table.columns[0].width =Inches(2.0) # set column widths
table.columns[1].width = Inches(4.0)
# 표 항목의 제목 작성 table.cell(0,0).text = '회사명' table.cell(0,1).text = '주식가격' # 쎌 내용 작성
table.cell(1,0).text = '네이버'
table.cell(1,1).text = '2200'
table.cell(2,0).text = '테슬라'
table.cell(2,1).text = '3300'
prs.save('test1.pptx')
