#3. 주식 보고서 자동 작성과 전송
import pandas as pd
def get_stock_code():
    stock_code = pd.read_html("https://kind.krx.co.kr/corpgeneral/corpList.do?method=download", header=0)[0]
    stock_code = stock_code[['회사명','종목코드']]
    # 영어로 변환
    stock_code = stock_code.rename(columns={'회사명':'company','종목코드':'code'})
    stock_code.code = stock_code.code.map('{:06d}'.format) # 6자리수로 통일 return stock_code
    return stock_code

# (company code) 내용 출력
code = get_stock_code()
print(code)

#######################################################################
#일별 시세 가져오기(특정 회사 종목 코드 이용)
import requests

agent_key = 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/114.0.0.0 Safari/537.36 Edg/114.0.1823.67'
def get_stock(code):
    df = pd.DataFrame()
    for page in range(1,21):
        url = 'http://finance.naver.com/item/sise_day.nhn?code={code}'.format(code=code)
        url ='{url}&page={page}'.format(url=url, page=page)
        print(url)
        header = {'User-Agent':'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/114.0.0.0 Safari/537.36 Edg/114.0.1823.67'} # 아래의 user-agent 설명 참조
        res = requests.get(url, headers=header)
        current_df = pd.read_html(res.text, header=0)[0]
        df = df._append(current_df, ignore_index = True)
        #df =df.append(current_df, ignore_index=True)
    return df

code = '005930' # 예) 삼성전자
df = get_stock(code)
print(df.head())

#######################################################################
#종목 코드와 일별 시세 정보 수집
def clean_data(df):
    df = df.dropna()
    df = df.rename(columns = {'날짜':'date', '종가':'close', '전일비':'diff', '시가':'open', '고가':'high',
                              '저가':'low', '거래량':'volume'})
    df[['close','diff','open','high','low','volume']] = df[['close','diff','open','high','low','volume']].astype(int)
    df['date'] =pd.to_datetime(df['date'])
    df = df.sort_values(by=['date'], ascending=True)
    return df
# 실습
company = '삼성전자'
stock_code = get_stock_code()
code = stock_code[stock_code.company==company].code.values[0].strip()
df = get_stock(code)
df = clean_data(df)
print(df)

#######################################################################
# 보고자료 (차트와 테이블) 준비
import matplotlib.pyplot as plt
from pandas.plotting import table
import os

# 차트 그리기
plt.figure(figsize=(10,4))
plt.plot(df['date'], df['close'])
plt.xlabel('date')
plt.ylabel('close')
# 차트 저장과 출력 (아래 경로 (C://images) 미리 만들어 두기)
chart_fname = os.path.join("C:/Users/LUIS/WorkSpace/AI_MBA/Artificial intelligence application practice", '{company}_chart.png'.format(company=company))
plt.savefig(chart_fname)
plt.show()
# 일별 시세 테이블 작성
plt.figure(figsize=(15,4))
ax = plt.subplot(111, frame_on=False)
ax.xaxis.set_visible(False)
ax.yaxis.set_visible(False)
df = df.sort_values(by=['date'], ascending=False)
table(ax, df.head(10), loc='center', cellLoc='center', rowLoc='center') # 일별 시세 테이블 저장
table_fname = os.path.join('C:/Users/LUIS/WorkSpace/AI_MBA/Artificial intelligence application practice', '{company}_table.png'.format(company=company))
plt.savefig(table_fname)

#######################################################################
#파워포인트로 보고서 작성
import datetime
from pptx import Presentation
from pptx.util import Inches
import os
########## 제목 슬라이드
today = datetime.datetime.today().strftime('%Y%m%d')
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]    # 제목 슬라이드
slide = prs.slides.add_slide(title_slide_layout) # 슬라이드 추가
# 제목 작성
title = slide.shapes.title
title.text = '주식 분석 보고서'
# 부제목 작성
subtitle = slide.placeholders[1]
subtitle.text = '작성일자: {date}'.format(date=today)

########## 차트와 테이블 슬라이드 추가
title_only_slide_layout = prs.slide_layouts[5]    # 제목 슬라이드
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes
shapes.title.text = '{company}, {close}원에 거래 마감'.format(company=company, close=df.iloc[0]['close'])
print(shapes.title.text)
# 차트 추가
left = Inches(0.5)
height = Inches(2.5)
width = Inches(9)
top = Inches(2)
pic = slide.shapes.add_picture(chart_fname, left, top, width=width, height=height)
# 테이블 추가
left = Inches(-1)
height = Inches(3)
width =Inches(12)
top = Inches(4)
pic = slide.shapes.add_picture(table_fname, left, top, width=width, height=height)
cursor_sp = slide.shapes[0]._element
cursor_sp.addprevious(pic._element)
####### 보고서 저장 (C://images 경로는 미리 존재하여야 함)
ppt_fname = os.path.join('C:/Users/LUIS/WorkSpace/AI_MBA/Artificial intelligence application practice', 'stock_report.pptx')
prs.save(ppt_fname)

#######################################################################
import os
from email.mime.multipart import MIMEMultipart
from email import encoders
from email.mime.text import MIMEText
from email.mime.image import MIMEImage
from email.mime.audio import MIMEAudio
from email.mime.base import MIMEBase
import smtplib

smtp_info = dict({"smtp_server":"smtp.naver.com",
                  "smtp_user_id":"luisfynn@naver.com",
                  "smtp_user_pw":"@kj0224kj@",
                  "smtp_port":587})

def send_email(smtp_info, msg):
    with smtplib.SMTP(smtp_info["smtp_server"], smtp_info["smtp_port"]) as server:
        server.starttls() #TLS 보안 연결
        server.login(smtp_info["smtp_user_id"], smtp_info["smtp_user_pw"]) # 로그인
        # 로그인된 서버에 이메일 전송
        response = server.sendmail(msg['from'], msg['to'], msg.as_string())

        if not response:
            print('이메일을 성공적으로 보냈습니다.')
        else:
            print(response)

def make_multimsg(msg_dict):
    multi = MIMEMultipart(_subtype='mixed')
    for key, value in msg_dict.items():
        if key == 'text':
            with open(value['filename'], encoding='utf-8') as fp:
                msg = MIMEText(fp.read(), _subtype=value['subtype'])
        elif key == 'image':
            with open(value['filename'], 'rb') as fp:
                msg = MIMEImage(fp.read(), _subtype=value['subtype'])
        elif key == 'audio':
            with open(value['filename'], 'rb') as fp:
                msg = MIMEAudio(fp.read(), _subtype=value['subtype'])
        else:
            with open(value['filename'], 'rb') as fp:
                msg = MIMEBase(value['maintype'], _subtype=value['subtype'])
                msg.set_payload(fp.read())
                encoders.encode_base64(msg)
        msg.add_header('Content-Disposition', 'attachment', filename=os.path.basename(value['filename']))
        multi.attach(msg)
    return multi


#보고서를 이메일로 전송
# 앞서 3개 함수 호출 필요: smtp_info, make_multimsg, send_email
title = '({date}).주식 분석 보고서'.format(date=today)
content = '주식 분석 보고서 입니다.'
sender = smtp_info['smtp_user_id']
#receiver ='myungkim@sogang.ac.kr'
receiver ='luisfynn1@gmail.com'
msg = MIMEText(_text = content, _charset = 'utf-8')
# 첨부파일 지정
msg_dict = {
'application': {'maintype': 'application', 'subtype': 'octect-stream', 'filename':'stock_report.pdf'}
}

# 첨부파일 추가
msg_dict['application']['filename'] = ppt_fname # PPT 형식 파일 지정
multi = make_multimsg(msg_dict) # msg_dict 안의 파일을 첨부함
multi['Subject'] = title
multi['From'] = sender
multi['To'] = receiver
multi.attach(msg)
# 이메일 전송
send_email(smtp_info, multi)